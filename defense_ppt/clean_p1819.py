# -*- coding: utf-8 -*-
"""Clean placeholders on P18/P19, ensure blank page truly blank."""
import sys
from pptx import Presentation

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
slides = list(prs.slides)

for idx in [17, 18]:  # P18 谢谢大家, P19 空白
    slide = slides[idx]
    removed = []
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape._element
            sp.getparent().remove(sp)
            removed.append(shape.name)
    print(f"Slide {idx+1}: removed placeholders {removed}, remaining shapes={len(slide.shapes)}")

prs.save(path)
print("saved clean")
