# -*- coding: utf-8 -*-
"""Round 4: P3 fill bottom, P4 fix table overlap+wrap, P5 stretch trace chain, P6 spread criteria."""
import sys
from pptx import Presentation
from pptx.util import Cm, Pt, Emu

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
slides = list(prs.slides)
sys.stdout.reconfigure(encoding='utf-8')

def cm(v): return Emu(int(v * 360000))

# ============ P3 (index 2): enlarge top cards, push timeline down, fill bottom ============
s = slides[2]
def move(shape, top=None, height=None, left=None, width=None):
    if top is not None: shape.top = cm(top)
    if height is not None: shape.height = cm(height)
    if left is not None: shape.left = cm(left)
    if width is not None: shape.width = cm(width)

def set_size(shape, pt):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                r.font.size = Pt(pt)

# stage title 16->18, body 13->14
set_size(s.shapes[5], 18); set_size(s.shapes[7], 18)   # stage titles
set_size(s.shapes[9], 14); set_size(s.shapes[11], 14); set_size(s.shapes[13], 14)  # body texts
# stage cards: move down + grow
for bg_title, bg_body, txt in [(6,8,9),(10,12,13)]:
    move(s.shapes[bg_title], top=4.0, height=1.3)
    move(s.shapes[bg_body], top=5.3, height=6.6)
    move(s.shapes[txt], top=5.55, height=6.35)
# widen body text frames a touch & add paragraph spacing
for txt in (9, 13):
    tf = s.shapes[txt].text_frame
    for p in tf.paragraphs:
        p.space_after = Pt(20)
# timeline: dates 11.22->12.0; dots 12.35->13.5; line 12.54->13.7; node desc 13.05->14.2
for i in (19,21,23,25,27,29):
    move(s.shapes[i], top=12.0, height=0.9)
    set_size(s.shapes[i], 14)
for i in (13,14,15,16,17,18):
    move(s.shapes[i], top=13.5)
move(s.shapes[12], top=13.7)  # connector line
for i in (20,22,24,26,28,30):
    move(s.shapes[i], top=14.2, height=1.7)
    set_size(s.shapes[i], 14)
move(s.shapes[31], top=18.1)  # bottom line
print("P3 done")

# ============ P4 (index 3): widen table, remove wrap, fix overlap ============
s = slides[3]
tbl = s.shapes[7].table
# widen: left 21.03, width 10.87 -> 12.0, cols [2.4, 9.6]
move(s.shapes[7], width=12.0)
tbl.columns[0].width = cm(2.4)
tbl.columns[1].width = cm(9.6)
# row heights for single-line (13pt): [0.95,1.0,1.0,1.0,1.0,0.95]
for row, h in zip(tbl.rows, [0.95,1.0,1.0,1.0,1.0,0.95]):
    row.height = cm(h)
# also right-aligned title/sub blocks keep position (no overlap now: table bottom=10.7)
# ensure selection-column runs that are 13pt fit: reduce to 12pt in value column to guarantee single line
for row in tbl.rows:
    for c in row.cells:
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and round(r.font.size.pt) == 13:
                    r.font.size = Pt(12)
print("P4 table bottom now:", round((s.shapes[7].top + sum(row.height for row in tbl.rows)) / 360000, 2))

# ============ P5 (index 4): stretch trace chain vertical ============
s = slides[4]
tops = [4.94, 7.4, 9.86, 12.32, 14.78, 17.24]
# frames are pairs: (bg, text) with same top; indexes 6/7,9/10,12/13,15/16,18/19,21/22
pairs = [(6,7),(9,10),(12,13),(15,16),(18,19),(21,22)]
for (bg, tx), t in zip(pairs, tops):
    move(s.shapes[bg], top=t)
    move(s.shapes[tx], top=t)
# connectors between frames: midpoint of gap
conn_tops = [6.7, 9.16, 11.62, 14.08, 16.54]
for idx, t in zip((8,11,14,17,20), conn_tops):
    move(s.shapes[idx], top=t)
print("P5 done")

# ============ P6 (index 5): spread four criteria, fill page ============
s = slides[5]
move(s.shapes[8], height=12.9)  # content box 5.01 -> 17.91
tf = s.shapes[8].text_frame
for p in tf.paragraphs:
    p.space_after = Pt(22)
move(s.shapes[9], top=18.0)  # bottom connector
print("P6 done")

prs.save(path)
print("saved")
