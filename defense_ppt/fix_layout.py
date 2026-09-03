# -*- coding: utf-8 -*-
"""Round 3-fix: repair overflow/overlap after font enlargement."""
import sys
from pptx import Presentation
from pptx.util import Cm, Pt

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
slides = list(prs.slides)
sys.stdout.reconfigure(encoding='utf-8')

def set_rows(tbl, heights_cm):
    for row, h in zip(tbl.rows, heights_cm):
        row.height = Cm(h)

# P2 table (shape 18): reduce row heights so bottom <= 18.9cm
t = slides[1].shapes[18].table
set_rows(t, [1.12, 1.33, 1.33, 1.33, 1.33, 1.33, 1.33, 1.33, 1.33])
print("P2 rows:", [round(r.height/360000,2) for r in t.rows], "total", round(sum(r.height for r in t.rows)/360000,2))

# P14 table (shape 6): reduce so bottom < 15.7
t = slides[13].shapes[6].table
set_rows(t, [1.30, 1.60, 1.60, 1.60, 2.00, 2.00])
print("P14 rows:", [round(r.height/360000,2) for r in t.rows], "total", round(sum(r.height for r in t.rows)/360000,2))

# P15 table (shape 22): reduce so bottom < 15.7
t = slides[14].shapes[22].table
set_rows(t, [1.22, 1.68, 1.68, 1.68, 1.68, 1.72])
print("P15 rows:", [round(r.height/360000,2) for r in t.rows], "total", round(sum(r.height for r in t.rows)/360000,2))

# P8 SQL code block: reduce font 10pt -> 9pt in code shapes to stop overflow
for sh in slides[7].shapes:
    if sh.has_text_frame:
        txt = sh.text_frame.text
        if 'REVOKE' in txt or 'GRANT' in txt or 'travelmate_traffic_app' in txt:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size and round(r.font.size.pt) == 10:
                        r.font.size = Pt(9)
            print("P8 code shape reduced:", sh.name, sh.text_frame.text[:30])

prs.save(path)
print("saved")
