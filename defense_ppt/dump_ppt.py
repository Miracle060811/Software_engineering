# -*- coding: utf-8 -*-
"""Dump detailed pptx structure: per-slide shapes with position, text, font."""
import sys
from pptx import Presentation
from pptx.util import Emu

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组.pptx"
out = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\.lark-slides\dump.txt"

prs = Presentation(path)
SLIDE_W = prs.slide_width  # EMU
SLIDE_H = prs.slide_height

def emu2cm(v):
    return round(Emu(v).cm, 2) if v is not None else None

lines = []
lines.append(f"SLIDE SIZE: {emu2cm(SLIDE_W)}x{emu2cm(SLIDE_H)} cm; total slides={len(prs.slides)}")
for i, slide in enumerate(prs.slides, 1):
    lines.append(f"\n========== SLIDE {i} ==========")
    for si, shape in enumerate(slide.shapes):
        info = f"[{si}] type={shape.shape_type} name='{shape.name}'"
        try:
            info += f" pos=({emu2cm(shape.left)},{emu2cm(shape.top)}) size=({emu2cm(shape.width)}x{emu2cm(shape.height)})"
        except Exception:
            pass
        lines.append(info)
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                runs_info = []
                for r in para.runs:
                    fs = r.font.size.pt if r.font.size else None
                    bold = r.font.bold
                    runs_info.append(f"'{r.text}'[sz={fs},b={bold}]")
                if runs_info:
                    lines.append(f"    p{pi}(align={para.alignment}): {' + '.join(runs_info)}")
        if shape.has_table:
            tbl = shape.table
            lines.append(f"    TABLE {len(tbl.rows)}x{len(tbl.columns)}")
            for ri, row in enumerate(tbl.rows):
                cells = []
                for ci, cell in enumerate(row.cells):
                    cells.append(cell.text.replace("\n", "⏎"))
                lines.append(f"      r{ri}: " + " | ".join(cells))
        if shape.shape_type == 13:  # picture
            lines.append(f"    PICTURE: {shape.image.filename if hasattr(shape.image,'filename') else 'embed'}")
        if shape.has_chart:
            lines.append("    CHART present")

with open(out, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print("written", out, "lines=", len(lines))
