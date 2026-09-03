# -*- coding: utf-8 -*-
"""Round 5: P3 timeline fill, P12 checklist spacing, P13 table unwrap, P14 bottom bar, P18 background."""
import sys
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
slides = list(prs.slides)
sys.stdout.reconfigure(encoding='utf-8')

def cm(v): return Cm(v)
def set_size(shape, pt):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                r.font.size = Pt(pt)

# ================= P3 (idx2): rework timeline, fill bottom =================
s = slides[2]
# dates 15pt & down
for i in (19,21,23,25,27,29):
    s.shapes[i].top = cm(12.3); s.shapes[i].height = cm(0.9)
    set_size(s.shapes[i], 15)
# dots & timeline line down
for i in (13,14,15,16,17,18):
    s.shapes[i].top = cm(13.7)
s.shapes[12].top = cm(13.9)
# node descriptions: bigger box, 16pt, vertically centered
for i in (20,22,24,26,28,30):
    s.shapes[i].top = cm(14.5); s.shapes[i].height = cm(2.6)
    set_size(s.shapes[i], 16)
    try:
        s.shapes[i].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
s.shapes[31].top = cm(18.6)
print("P3 done")

# ================= P12 (idx11): enlarge checklist, spread lines =================
s = slides[11]
box = s.shapes[8]  # AutoShape 10: 6-bullet checklist
set_size(box, 14)
tf = box.text_frame
for p in tf.paragraphs:
    p.space_after = Pt(14)
print("P12 done")

# ================= P13 (idx12): widen tool column to avoid wrap =================
s = slides[12]
tbl = None
for sh in s.shapes:
    if sh.has_table:
        tbl = sh.table; break
if tbl:
    widths = [6.5, 8.8, 3.9, 3.9, 7.2]
    for col, w in zip(tbl.columns, widths):
        col.width = cm(w)
    # drop long tool texts to 12pt for safety
    for ri in (5,6):
        c = tbl.rows[ri].cells[1]
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and round(r.font.size.pt) >= 13:
                    r.font.size = Pt(12)
    print("P13 cols:", [round(Emu(c.width).cm,2) for c in tbl.columns], "total", round(sum(Emu(c.width) for c in tbl.columns)/360000,2))

# ================= P14 (idx13): bigger bottom bar, comfortable text =================
s = slides[13]
s.shapes[21].height = cm(2.9)        # bar 15.52 -> 18.42
s.shapes[22].top = cm(15.95); s.shapes[22].height = cm(2.0)   # label
set_size(s.shapes[22], 16)
s.shapes[23].top = cm(15.95); s.shapes[23].height = cm(2.0)   # desc text
set_size(s.shapes[23], 14)
try:
    s.shapes[23].text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
except Exception:
    pass
print("P14 done")

# ================= P18 (idx17): background to F2F4F8 like others =================
s = slides[17]
s.background.fill.solid()
s.background.fill.fore_color.rgb = RGBColor(0xF2, 0xF4, 0xF8)
print("P18 bg set")

prs.save(path)
print("saved")
