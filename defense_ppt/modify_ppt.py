# -*- coding: utf-8 -*-
"""Modify TravelMate defense PPT: fix content per project docs, remove draft-only
phrases (待现场实测/评分标准 etc), add 谢谢大家 + blank page."""
import sys
from pptx import Presentation
from pptx.util import Pt, Emu

SRC = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组.pptx"
OUT = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"

prs = Presentation(SRC)
slides = list(prs.slides)

def para_set_text(para, new_text):
    """Replace paragraph text, keeping first run's formatting."""
    if para.runs:
        para.runs[0].text = new_text
        for r in para.runs[1:]:
            r.text = ""
    else:
        para.add_run().text = new_text

def shape_set_text(shape, lines, keep_empty=False):
    """lines: list of (para_index, text)."""
    tf = shape.text_frame
    for pi, text in lines:
        if pi < len(tf.paragraphs):
            para_set_text(tf.paragraphs[pi], text)

def cell_set_text(cell, new_text):
    """Set cell text with \n as paragraph breaks, preserving font from old first run."""
    tf = cell.text_frame
    ref = None
    for p in tf.paragraphs:
        if p.runs:
            ref = p.runs[0]
            break
    # clear paragraphs beyond first
    first = tf.paragraphs[0]
    # remove extra paragraphs
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    first.clear()
    lines = new_text.split("\n")
    for i, line in enumerate(lines):
        p = first if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if ref is not None:
            r.font.size = ref.font.size
            r.font.bold = ref.font.bold
            r.font.name = ref.font.name
            try:
                r.font.color.rgb = ref.font.color.rgb
            except Exception:
                pass

def get_table(slide, shape_idx):
    return slide.shapes[shape_idx].table

def set_table_cell(slide, shape_idx, row, col, text):
    tbl = get_table(slide, shape_idx)
    cell_set_text(tbl.rows[row].cells[col], text)

# ============ P2 ============
s = slides[1]
shape_set_text(s.shapes[2], [(0, "八项验收要求与证据总览")])
shape_set_text(s.shapes[9], [(0, "245")])
shape_set_text(s.shapes[11], [(0, "14")])
shape_set_text(s.shapes[13], [(0, "325")])
# Table 20 (shape 18): rows r2,r4,r5,r6,r8, col 1
set_table_cell(s, 18, 2, 1, "　14 个 Job 门禁链，测试失败即阻断合并与部署")
set_table_cell(s, 18, 4, 1, "　34 张业务表一表一主（另 1 张媒体表暂留单体）、6 套独立 MySQL，跨库访问被账号拒绝")
set_table_cell(s, 18, 5, 1, "　单元 / 集成 / E2E 三层共 325 项全绿，需求可追溯")
set_table_cell(s, 18, 6, 1, "　6 个 HPA：副本 2→6、CPU 阈值 60%，09-02 已实测通过")
set_table_cell(s, 18, 8, 1, "　同机同数据同脚本，2 场景 × 3 轮已实测出数；并发下单防超卖验证")

# ============ P3 ============
s = slides[2]
shape_set_text(s.shapes[2], [(0, "十天两阶段实施路线")])
shape_set_text(s.shapes[5], [(0, "阶段一 · 单体复活与工程化（08-22 → 08-27）")])
shape_set_text(s.shapes[9], [(0, "阶段二 · 微服务与云原生（08-28 → 09-01）")])
# timeline nodes: shape26 (08-31 text), shape28 (09-01 text)
shape_set_text(s.shapes[26], [(0, "六服务上线"), (1, "部署就绪")])
shape_set_text(s.shapes[28], [(0, "E2E 17/17"), (1, "边界验证")])
shape_set_text(s.shapes[32], [(0, "交付物按 01 需求 / 02 设计 / 03 代码 / 04 测试 / 05 管理 / 06 交付 六个目录归档；每日站会、看板流转、PR 评审与 AI 使用记录全程留痕。")])

# ============ P4 ============
s = slides[3]
shape_set_text(s.shapes[2], [(0, "总体架构与技术选型")])
set_table_cell(s, 7, 3, 1, "　MySQL 8（单体 8.0 / 微服务 8.4）/ Redis 7（缓存与锁）")

# ============ P5 ============
s = slides[4]
shape_set_text(s.shapes[2], [(0, "UC01—UC19 用例全景与追溯链")])
tbl = get_table(s, 4)
rows_data = [
    (1, "账号\n与权限", "UC01 注册登录与账户安全\nUC16 常用旅客　UC17 主页与关注"),
    (2, "交通\n出行", "UC02 航班查询与预订\nUC03 火车票与候补　UC04 交通订单支付退改"),
    (3, "当地\n服务", "UC05 酒店订房　UC06 酒店订单与库存回补\nUC07 景点购票　UC08 一日游/周边游"),
    (4, "权益\n与 AI", "UC10 优惠券领取核销\nUC11 AI 行程生成　UC12 AI 客服对话"),
    (5, "旅行\n社区", "UC09 评价与举报　UC14 游记发布审核\nUC15 点赞收藏评论"),
    (6, "后台\n与消息", "UC13 通知与私信　UC18 管理后台\nUC19 内容安全与可观测性"),
]
for r, dom, uc in rows_data:
    cell_set_text(tbl.rows[r].cells[0], dom)
    cell_set_text(tbl.rows[r].cells[1], uc)

# ============ P6/P7 ============
shape_set_text(slides[5].shapes[2], [(0, "按业务能力的六服务划分")])
shape_set_text(slides[6].shapes[2], [(0, "六服务职责、端口与表归属")])

# ============ P8 ============
s = slides[7]
shape_set_text(s.shapes[2], [(0, "数据所有权：一表一主与账号隔离")])
tbl8 = get_table(s, 18)
acct = [("travelmate_identity_app", "traffic.tm_flight"),
        ("travelmate_traffic_app", "local.tm_hotel"),
        ("travelmate_local_app", "ai.tm_notification"),
        ("travelmate_ai_app", "identity.tm_user")]
for i, (acc, tab) in enumerate(acct, start=1):
    cell_set_text(tbl8.rows[i].cells[0], "　" + acc)
    cell_set_text(tbl8.rows[i].cells[1], tab)
# code block shape 20 p3
shape_set_text(s.shapes[20], [(3, "　ON travelmate_traffic.* TO 'travelmate_traffic_app'@'%';")])

# ============ P9 ============
shape_set_text(slides[8].shapes[2], [(0, "跨服务协作与失败补偿契约")])

# ============ P10 ============
s = slides[9]
shape_set_text(s.shapes[2], [(0, "CI/CD 四道质量闸")])
shape_set_text(s.shapes[6], [(0, "14 个 Job · 四道闸")])

# ============ P11 ============
s = slides[10]
shape_set_text(s.shapes[2], [(0, "流水线实证：Run #217 全绿")])
shape_set_text(s.shapes[8], [(1, "全部 Job 通过，无跳过项")])

# ============ P12 ============
shape_set_text(slides[11].shapes[2], [(0, "Kubernetes 部署实证")])

# ============ P13 ============
s = slides[12]
shape_set_text(s.shapes[2], [(0, "六层自动化测试结果")])
set_table_cell(s, 4, 4, 2, "17")
set_table_cell(s, 4, 5, 2, "100")
shape_set_text(s.shapes[6], [(2, "UC02 越权访问被拦截；UC04 订单状态机完整流转、退改可回滚。")])

# ============ P14 ============
s = slides[13]
shape_set_text(s.shapes[2], [(0, "HPA 自动扩缩容实测")])
shape_set_text(s.shapes[7], [(0, "实测五步（脚本自动留存三态）")])
shape_set_text(s.shapes[22], [(0, "已实测通过（09-02）")])
shape_set_text(s.shapes[23], [(0, "09-02 已实测：k6 加压后副本由 2 扩至 6，停止加压后经 120s 稳定窗口缩回 2，CPU 峰值 237%，实验判定 passed=true；完整记录与截图存于 04_tests/stress/results。")])

# ============ P15 ============
shape_set_text(slides[14].shapes[2], [(0, "故障隔离实证：身份服务宕机")])

# ============ P16 ============
s = slides[15]
shape_set_text(s.shapes[2], [(0, "单体 vs 微服务性能对比")])
shape_set_text(s.shapes[8], [(0, "结果记录表（2026-09-02 实测，三轮中位数）")])
tbl16 = get_table(s, 10)
perf = [
    (1, "航班查询", "单体", "30.964", "20.533", "0%"),
    (2, "", "微服务", "31.565", "4.595", "0%"),
    (3, "酒店查询", "单体", "30.803", "22.794", "0%"),
    (4, "", "微服务", "31.564", "4.436", "0%"),
    (5, "并发下单", "单体", "——", "——", "——"),
    (6, "", "微服务", "——", "——", "——"),
]
for r, sc, ver, tps, p95, err in perf:
    cell_set_text(tbl16.rows[r].cells[0], sc)
    cell_set_text(tbl16.rows[r].cells[1], ver)
    cell_set_text(tbl16.rows[r].cells[2], tps)
    cell_set_text(tbl16.rows[r].cells[3], p95)
    cell_set_text(tbl16.rows[r].cells[4], err)

# ============ P17 ============
s = slides[16]
shape_set_text(s.shapes[2], [(0, "过程管理与演示路线")])
shape_set_text(s.shapes[6], [(1, "245 次提交，monolith-start 与 microservices-phase1 两个标签冻结阶段基线")])
shape_set_text(s.shapes[6], [(5, "性能对比与 HPA 扩缩容均已有 09-02 实测记录，结论可复现")])
shape_set_text(s.shapes[16], [(0, "重点用例：UC04 交通订单履约、UC11 AI 行程生成（约 1.5 分钟）")])

# ============ add 谢谢大家 page (18) & blank page (19) ============
# pick a blank-ish layout
blank = None
for ly in prs.slide_layouts:
    if ly.name and ("Blank" in ly.name or "空白" in ly.name):
        blank = ly
        break
if blank is None:
    blank = prs.slide_layouts[-1]

# Slide 18: 谢谢大家
s18 = prs.slides.add_slide(blank)
# white background
from pptx.enum.shapes import MSO_SHAPE
SW, SH = prs.slide_width, prs.slide_height

def add_text(slide, x, y, w, h, lines, align=None):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN
    for i, (txt, size, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if align is None else align
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Microsoft YaHei"
    return tb

add_text(s18, int(SW*0.15), int(SH*0.30), int(SW*0.70), int(SH*0.4), [
    ("谢谢大家", 66, True),
])

add_text(s18, int(SW*0.15), int(SH*0.55), int(SW*0.70), int(SH*0.3), [
    ("TravelMate：从单体到六服务的云原生改造", 22, False),
    ("第 5 组 · Miracle 小组　|　2026-09-04 最终答辩", 16, False),
])

# Slide 19: blank page
prs.slides.add_slide(blank)

prs.save(OUT)
print("saved:", OUT)
print("total slides now:", len(prs.slides))
