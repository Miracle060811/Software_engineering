# -*- coding: utf-8 -*-
"""Round 3: enlarge fonts globally, scale table row heights, then save.
Mapping keeps hierarchy: title/section/body/note all +30~50%."""
import sys
from pptx import Presentation
from pptx.util import Pt, Emu

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
sys.stdout.reconfigure(encoding='utf-8')

M = {22: 26, 13: 16, 11: 14, 10: 13, 9: 12, 8: 10,
     36: 40, 30: 34, 18: 20, 14: 17, 12: 15, 66: 80, 16: 19}

def ms(pt):
    return M.get(pt, pt + 3)

changed = 0
for si, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        ns = ms(round(r.font.size.pt))
                        if ns != round(r.font.size.pt):
                            r.font.size = Pt(ns); changed += 1
        if sh.has_table:
            tbl = sh.table
            for row in tbl.rows:
                row.height = int(row.height * 1.18)
                for c in row.cells:
                    for p in c.text_frame.paragraphs:
                        for r in p.runs:
                            if r.font.size:
                                ns = ms(round(r.font.size.pt))
                                if ns != round(r.font.size.pt):
                                    r.font.size = Pt(ns); changed += 1
print("font size changes:", changed)
prs.save(path)
print("saved")
