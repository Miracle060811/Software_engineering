# -*- coding: utf-8 -*-
"""Round 4-fix: correct P3 layout (previous round had wrong shape indices)."""
import sys
from pptx import Presentation
from pptx.util import Cm, Pt

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
s = list(prs.slides)[2]
sys.stdout.reconfigure(encoding='utf-8')

def cm(v): return Cm(v)
def set_size(shape, pt):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.size:
                r.font.size = Pt(pt)

# ---- correct shape indices on slide 3 ----
# 4=AutoShape6 stage1 title bg, 5=stage1 title txt, 6=stage1 body bg, 7=stage1 body txt
# 8=stage2 title bg, 9=stage2 title txt, 10=stage2 body bg, 11=stage2 body txt
# 12=timeline line, 13..18=dots, 19/21/23/25/27/29=dates, 20/22/24/26/28/30=nodes, 31=bottom line

# stage cards: titles enlarged, body moved down+grew, widened
for ti_bg, ti_tx, bo_bg, bo_tx in [(4,5,6,7),(8,9,10,11)]:
    s.shapes[ti_bg].top = cm(4.0); s.shapes[ti_bg].height = cm(1.3)
    s.shapes[ti_tx].top = cm(4.0); s.shapes[ti_tx].height = cm(1.3)
    s.shapes[bo_bg].top = cm(5.3); s.shapes[bo_bg].height = cm(6.6)
    s.shapes[bo_tx].top = cm(5.55); s.shapes[bo_tx].height = cm(6.35)
    s.shapes[bo_tx].width = cm(14.2)  # widen to avoid wrap at 14pt
    set_size(s.shapes[ti_tx], 18)
    set_size(s.shapes[bo_tx], 14)
    tf = s.shapes[bo_tx].text_frame
    for p in tf.paragraphs:
        p.space_after = Pt(18)

# timeline down
s.shapes[12].top = cm(13.7)  # timeline line
for i in (13,14,15,16,17,18):
    s.shapes[i].top = cm(13.5)  # dots
for i in (19,21,23,25,27,29):
    s.shapes[i].top = cm(12.0); s.shapes[i].height = cm(0.9)
    set_size(s.shapes[i], 14)
for i in (20,22,24,26,28,30):
    s.shapes[i].top = cm(14.2); s.shapes[i].height = cm(1.7)
    set_size(s.shapes[i], 14)
s.shapes[31].top = cm(18.1)  # bottom line

prs.save(path)
print("P3 fixed and saved")
