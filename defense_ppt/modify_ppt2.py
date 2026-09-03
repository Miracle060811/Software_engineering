# -*- coding: utf-8 -*-
"""Round 2: delete footer texts per user request, fix P5 table layout, remove watermark."""
import sys
from pptx import Presentation
from pptx.util import Emu, Pt

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
slides = list(prs.slides)
sys.stdout.reconfigure(encoding='utf-8')

def delete_shape_by_name(slide, name):
    for shape in list(slide.shapes):
        if shape.name == name:
            sp = shape._element
            sp.getparent().remove(sp)
            return True
    return False

def set_cell_singleline(cell, text):
    """Set cell to single line, keep font of first run."""
    tf = cell.text_frame
    ref = None
    for p in tf.paragraphs:
        if p.runs:
            ref = p.runs[0]
            break
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    first.clear()
    r = first.add_run()
    r.text = text
    if ref is not None:
        r.font.size = ref.font.size
        r.font.bold = ref.font.bold
        r.font.name = ref.font.name
        try:
            r.font.color.rgb = ref.font.color.rgb
        except Exception:
            pass
    # vertical center
    try:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

# 1) P3: delete bottom delivery-archive note (AutoShape 34)
ok3 = delete_shape_by_name(slides[2], "AutoShape 34")
print("P3 delete bottom note:", ok3)

# 2) P5: delete matrix-file note (AutoShape 25) + fix table col-0 to single-line domain names
ok5 = delete_shape_by_name(slides[4], "AutoShape 25")
print("P5 delete matrix note:", ok5)
tbl = slides[4].shapes[4].table
domains = ["账号与权限", "交通出行", "当地服务", "权益与 AI", "旅行社区", "后台与消息"]
for i, d in enumerate(domains, start=1):
    set_cell_singleline(tbl.rows[i].cells[0], d)
    print("P5 r%d col0:" % i, tbl.rows[i].cells[0].text_frame.text)

# 3) P6: delete bottom note (AutoShape 12)
ok6 = delete_shape_by_name(slides[5], "AutoShape 12")
print("P6 delete bottom note:", ok6)

# 4) P7: delete bottom note (AutoShape 7)
ok7 = delete_shape_by_name(slides[6], "AutoShape 7")
print("P7 delete bottom note:", ok7)

# 5) P17: delete watermark picture (Picture 27)
ok17 = delete_shape_by_name(slides[16], "Picture 27")
print("P17 delete watermark:", ok17)

prs.save(path)
print("saved")
