# -*- coding: utf-8 -*-
"""Verify modified pptx: dump titles and key changed texts."""
import sys
from pptx import Presentation

path = r"D:\Campus\grade2\26Spring\Software_engineering\defense_ppt\export\TravelMate云原生改造-最终答辩-第5组_已修改.pptx"
prs = Presentation(path)
sys.stdout.reconfigure(encoding='utf-8')

def cell_text(cell):
    return " / ".join(p.text for p in cell.text_frame.paragraphs if p.text)

print("TOTAL SLIDES:", len(prs.slides))
for i, slide in enumerate(prs.slides, 1):
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            # title is usually the biggest bold text near top
            t = shape.text_frame.text.strip().replace("\n", " ")
            if title is None and len(t) < 40 and shape.top is not None and shape.top < 2000000:
                title = t
                break
    print(f"--- Slide {i} title-hint: {title}")

print("\n=== key checks ===")
def find_slide(idx):
    return list(prs.slides)[idx - 1]

# P2 numbers
s = find_slide(2)
nums = []
for si in [9, 11, 13]:
    nums.append(s.shapes[si].text_frame.text)
print("P2 numbers (245/14/325 expected):", nums)
# P2 table row2/r4/r6/r8
t = s.shapes[18].table
print("P2 r2:", cell_text(t.rows[2].cells[1])[:50])
print("P2 r4:", cell_text(t.rows[4].cells[1])[:60])
print("P2 r6:", cell_text(t.rows[6].cells[1])[:50])
print("P2 r8:", cell_text(t.rows[8].cells[1])[:60])

# P3
s = find_slide(3)
print("P3 stage1:", s.shapes[5].text_frame.text)
print("P3 stage2:", s.shapes[9].text_frame.text)
print("P3 node0831:", " / ".join(p.text for p in s.shapes[26].text_frame.paragraphs))
print("P3 node0901:", " / ".join(p.text for p in s.shapes[28].text_frame.paragraphs))
print("P3 bottom:", s.shapes[32].text_frame.text[:80])

# P5 table
s = find_slide(5)
t = s.shapes[4].table
for r in range(1, 7):
    print(f"P5 r{r}:", cell_text(t.rows[r].cells[0]), "=>", cell_text(t.rows[r].cells[1])[:80])

# P8
s = find_slide(8)
t = s.shapes[18].table
for r in range(1, 5):
    print("P8 r%d:" % r, cell_text(t.rows[r].cells[0]), "|", cell_text(t.rows[r].cells[1]))

# P13
s = find_slide(13)
t = s.shapes[4].table
print("P13 r4 c2 (expect 17):", cell_text(t.rows[4].cells[2]))
print("P13 r5 c2 (expect 100):", cell_text(t.rows[5].cells[2]))
print("P13 bottom p2:", s.shapes[6].text_frame.paragraphs[2].text)

# P14
s = find_slide(14)
print("P14 badge:", s.shapes[22].text_frame.text)
print("P14 desc:", s.shapes[23].text_frame.text[:80])

# P16
s = find_slide(16)
t = s.shapes[10].table
for r in range(1, 7):
    print("P16 r%d:" % r, " | ".join(cell_text(t.rows[r].cells[c]) for c in range(5)))
print("P16 head:", s.shapes[8].text_frame.text)

# P17
s = find_slide(17)
print("P17 p1:", s.shapes[6].text_frame.paragraphs[1].text)
print("P17 p5:", s.shapes[6].text_frame.paragraphs[5].text)
print("P17 step4:", s.shapes[16].text_frame.text)

# P18/P19
s18 = find_slide(18)
print("P18 texts:", [sh.text_frame.text for sh in s18.shapes if sh.has_text_frame and sh.text_frame.text.strip()])
s19 = find_slide(19)
print("P19 shape count:", len(s19.shapes))
